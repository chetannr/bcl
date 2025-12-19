#!/usr/bin/env python3
"""
BCL Re-Auction 2025 - Generate Player Slides
Creates slides for all players using the first slide as a template.
Replaces the left-hand side player photo and updates text content.
"""

import json
import sys
import copy
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from PIL import Image


def load_player_data(json_path: Path) -> List[Dict]:
    """Load player data from JSON file."""
    try:
        with open(json_path, 'r', encoding='utf-8') as f:
            players = json.load(f)
        print(f"Loaded {len(players)} players from {json_path.name}")
        return players
    except Exception as e:
        print(f"Error loading player data: {e}")
        sys.exit(1)


def scan_player_images(images_dir: Path) -> Dict[str, Path]:
    """Scan directory for player images and create mobile number mapping."""
    image_map = {}
    image_extensions = {'.jpg', '.jpeg', '.png', '.JPG', '.JPEG', '.PNG'}
    
    for image_file in images_dir.glob('*'):
        if image_file.suffix in image_extensions:
            # Extract mobile number from filename (e.g., "6360452535.jpg" -> "6360452535")
            mobile = image_file.stem
            if mobile.isdigit():
                image_map[mobile] = image_file
    
    print(f"Found {len(image_map)} player images")
    return image_map


def find_first_image_shape(slide) -> Optional:
    """Find the leftmost image shape in the slide (left-hand side player photo)."""
    image_shapes = []
    
    # Collect all image shapes with their positions
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_shapes.append((shape.left, shape))
        # Also check grouped shapes
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_shapes.append((sub_shape.left, sub_shape))
    
    if not image_shapes:
        return None
    
    # Return the leftmost image (smallest left position)
    # This should be the player photo on the left-hand side
    leftmost_image = min(image_shapes, key=lambda x: x[0])
    return leftmost_image[1]


def replace_image_in_shape(shape, image_path: Path, slide) -> bool:
    """Replace the image in a shape with a new image file, maintaining natural aspect ratio."""
    try:
        # Check if image format is supported by python-pptx
        ext = image_path.suffix.lower()
        supported_formats = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif', '.wmf'}
        if ext not in supported_formats:
            print(f"  Warning: Unsupported image format: {ext} for {image_path.name}")
            return False
        
        # Verify image file exists and is readable
        if not image_path.exists():
            print(f"  Warning: Image file does not exist: {image_path}")
            return False
        
        # Save shape properties before removal (available space)
        available_left = shape.left
        available_top = shape.top
        available_width = shape.width
        available_height = shape.height
        
        # Get the actual image dimensions to calculate aspect ratio
        try:
            with Image.open(image_path) as img:
                img_width, img_height = img.size
                img_aspect_ratio = img_width / img_height
        except Exception as e:
            print(f"  Warning: Could not read image dimensions for {image_path.name}: {e}")
            # Fall back to using available space dimensions
            img_aspect_ratio = available_width / available_height
        
        # Calculate available space aspect ratio
        available_aspect_ratio = available_width / available_height
        
        # Calculate new dimensions maintaining natural aspect ratio
        if img_aspect_ratio > available_aspect_ratio:
            # Image is wider - fit to width
            new_width = available_width
            new_height = available_width / img_aspect_ratio
        else:
            # Image is taller - fit to height
            new_height = available_height
            new_width = available_height * img_aspect_ratio
        
        # Center the image within the available space
        new_left = available_left + (available_width - new_width) // 2
        new_top = available_top + (available_height - new_height) // 2
        
        # Get the parent element before removal
        parent = shape._element.getparent()
        
        # Remove the old shape element from the tree
        parent.remove(shape._element)
        
        # Add new image at calculated position with natural aspect ratio
        # This creates a new relationship automatically and embeds the image
        new_picture = slide.shapes.add_picture(str(image_path), new_left, new_top, new_width, new_height)
        
        # Verify the picture was actually added and has a valid image
        if new_picture is None:
            return False
        
        # Verify the image relationship exists
        try:
            test_img = new_picture.image
            if test_img is None or len(test_img.blob) == 0:
                return False
        except Exception:
            return False
        
        return True
        
    except Exception as e:
        print(f"  Warning: Could not replace image for {image_path.name}: {e}")
        import traceback
        traceback.print_exc()
        return False


def create_player_info_table(slide, player: Dict, presentation: Presentation) -> None:
    """Create a table at the bottom of the slide with player information."""
    player_name = player.get('Name', '')
    age = player.get('Age', '')
    category = player.get('Category', '')
    phone = player.get('Ph', '')
    
    # Remove any existing tables at the bottom (if regenerating)
    shapes_to_remove = []
    for shape in slide.shapes:
        if hasattr(shape, 'has_table') and shape.has_table:
            # Check if it's at the bottom (likely our player info table)
            if shape.top > Inches(6):  # Near bottom of slide
                shapes_to_remove.append(shape)
    
    for shape in shapes_to_remove:
        sp = shape._element
        sp.getparent().remove(sp)
    
    # Standard PowerPoint slide dimensions (16:9 aspect ratio)
    # Width: 10 inches, Height: 7.5 inches
    slide_width = Inches(10)
    slide_height = Inches(7.5)
    
    # Table dimensions: 100% width, positioned higher up on slide
    table_width = slide_width
    table_height = Inches(1.2)  # Adjust height as needed
    table_left = Inches(0)
    table_top = Inches(4.5)  # Position at 4.5 inches from top (higher up, more visible)
    
    # Create table with 2 rows and 2 columns
    rows = 2
    cols = 2
    
    # Add table to slide
    table_shape = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
    table = table_shape.table
    
    # Set column widths to 50% each (must be integer)
    col_width = int(table_width / 2)
    table.columns[0].width = col_width
    table.columns[1].width = col_width
    
    # Fill in the data with actual player information
    # Row 1: Name | Age
    table.cell(0, 0).text = player_name if player_name else ''
    table.cell(0, 1).text = age if age else ''
    
    # Row 2: Category | Phone
    table.cell(1, 0).text = category if category else ''
    table.cell(1, 1).text = phone if phone else ''
    
    # Format cells - center align text and make visible
    for row_idx in range(rows):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            # Get the text that was just set
            cell_text = cell.text.strip() if cell.text else ''
            
            # Ensure text is set - clear and recreate
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            
            # Add text with formatting
            run = p.add_run()
            run.text = cell_text
            run.font.size = Pt(24)  # Larger font for visibility
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 0, 0)  # Black text
            
            # Set vertical alignment to middle
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            # Set cell fill to white for better visibility
            try:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
            except:
                pass
    
    # Ensure table is visible by setting it to be on top layer
    # Tables added last should be on top, but we can also verify
    try:
        # Move table to end of shape tree (brings to front)
        table_element = table_shape._element
        parent = table_element.getparent()
        parent.remove(table_element)
        parent.append(table_element)
    except:
        pass


def duplicate_slide(presentation: Presentation, source_slide) -> 'Slide':
    """
    Duplicate a slide in the presentation.
    
    Args:
        presentation: Presentation object
        source_slide: Source slide to duplicate
        
    Returns:
        New duplicated slide
    """
    # Get slide layout from source
    slide_layout = source_slide.slide_layout
    
    # Create new slide with same layout
    new_slide = presentation.slides.add_slide(slide_layout)
    
    # Remove default shapes from the new slide
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)
    
    # Copy shapes from source slide
    # Note: Image relationships will be broken, but we'll replace images anyway
    for shape in source_slide.shapes:
        el = shape._element
        newel = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    
    return new_slide


def create_player_slide(presentation: Presentation, template_slide, player: Dict, image_map: Dict[str, Path]) -> Tuple[bool, str]:
    """
    Create a new slide for a player based on the template.
    
    Returns:
        Tuple of (success: bool, message: str)
    """
    try:
        # Duplicate the template slide
        new_slide = duplicate_slide(presentation, template_slide)
        
        # Now find and replace the first image
        first_image_shape = find_first_image_shape(new_slide)
        phone = player.get('Ph', '')
        
        if first_image_shape and phone in image_map:
            image_path = image_map[phone]
            if replace_image_in_shape(first_image_shape, image_path, new_slide):
                image_replaced = True
            else:
                image_replaced = False
        else:
            image_replaced = False
            if not first_image_shape:
                print(f"  Warning: No image shape found for {player.get('Name', 'Unknown')}")
            elif phone not in image_map:
                print(f"  Warning: Image not found for {player.get('Name', 'Unknown')} (Phone: {phone})")
        
        # Remove any existing text shapes that might contain old player data
        # (These are copied from the template and should be removed)
        shapes_to_remove = []
        for shape in list(new_slide.shapes):  # Use list() to avoid iteration issues
            if hasattr(shape, 'text') and shape.has_text_frame:
                text = shape.text.strip()
                # Remove text shapes that look like player info (contain names, ages, categories, phones)
                if text and (
                    any(keyword in text.lower() for keyword in ['name:', 'age:', 'category:', 'ph:', 'phone:']) or
                    len([c for c in text if c.isdigit()]) >= 10  # Contains phone number (10+ digits)
                ):
                    shapes_to_remove.append(shape)
        
        # Remove the identified shapes
        for shape in shapes_to_remove:
            try:
                sp = shape._element
                parent = sp.getparent()
                if parent is not None:
                    parent.remove(sp)
            except Exception as e:
                # Shape might already be removed, continue
                pass
        
        # Remove any existing tables (if regenerating)
        tables_to_remove = []
        for shape in list(new_slide.shapes):
            if hasattr(shape, 'has_table') and shape.has_table:
                if shape.top > Inches(5):  # Near bottom
                    tables_to_remove.append(shape)
        
        for shape in tables_to_remove:
            try:
                sp = shape._element
                parent = sp.getparent()
                if parent is not None:
                    parent.remove(sp)
            except:
                pass
        
        # Create player info table at the bottom (even if image replacement failed)
        # This is added last so it appears on top
        try:
            create_player_info_table(new_slide, player, presentation)
        except Exception as table_error:
            print(f"  Warning: Could not create table for {player.get('Name', 'Unknown')}: {table_error}")
        
        player_name = player.get('Name', 'Unknown')
        status = "✓" if image_replaced else "⚠ (no image)"
        return True, f"{status} {player_name} ({phone})"
        
    except Exception as e:
        player_name = player.get('Name', 'Unknown')
        error_msg = str(e)
        return False, f"✗ {player_name}: {error_msg}"


def main():
    """Main function to generate player slides."""
    # Set up paths
    script_dir = Path(__file__).parent
    pptx_path = script_dir / "Re-Auction-2025-BCL-Players.pptx"
    json_path = script_dir / "players_data.json"
    images_dir = script_dir
    
    # Validate input files
    if not pptx_path.exists():
        print(f"Error: PowerPoint file not found: {pptx_path}")
        sys.exit(1)
    
    if not json_path.exists():
        print(f"Error: Player data file not found: {json_path}")
        sys.exit(1)
    
    print("=" * 60)
    print("BCL Re-Auction 2025 - Generate Player Slides")
    print("=" * 60)
    print()
    
    # Load player data
    players = load_player_data(json_path)
    
    # Scan for player images
    image_map = scan_player_images(images_dir)
    print()
    
    # Load PowerPoint presentation
    print(f"Loading PowerPoint: {pptx_path.name}")
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"Error loading PowerPoint: {e}")
        sys.exit(1)
    
    # Get template slide (first slide)
    if len(prs.slides) == 0:
        print("Error: PowerPoint file has no slides")
        sys.exit(1)
    
    template_slide = prs.slides[0]
    print(f"Using slide 1 as template")
    print(f"Current slides in presentation: {len(prs.slides)}")
    print()
    
    # Statistics
    stats = {
        'total_players': len(players),
        'slides_created': 0,
        'images_found': 0,
        'images_missing': 0,
        'errors': 0
    }
    
    # Generate slides for each player
    print("Generating player slides...")
    print("-" * 60)
    
    for i, player in enumerate(players, start=1):
        phone = player.get('Ph', '')
        
        # Skip players without phone numbers
        if not phone or phone == '?' or not phone.isdigit():
            print(f"  [{i:3d}/{len(players)}] ⚠ Skipping {player.get('Name', 'Unknown')} - invalid phone number")
            stats['errors'] += 1
            continue
        
        has_image = phone in image_map
        
        try:
            success, message = create_player_slide(prs, template_slide, player, image_map)
            
            if success:
                stats['slides_created'] += 1
                if has_image:
                    stats['images_found'] += 1
                else:
                    stats['images_missing'] += 1
            else:
                stats['errors'] += 1
            
            print(f"  [{i:3d}/{len(players)}] {message}")
        except Exception as e:
            stats['errors'] += 1
            player_name = player.get('Name', 'Unknown')
            print(f"  [{i:3d}/{len(players)}] ✗ {player_name}: Error - {str(e)}")
    
    print("-" * 60)
    print()
    
    # Save the updated presentation
    output_path = pptx_path  # Overwrite original
    print(f"Saving updated presentation to: {output_path.name}")
    try:
        prs.save(output_path)
        print("✓ Presentation saved successfully")
    except Exception as e:
        print(f"✗ Error saving presentation: {e}")
        sys.exit(1)
    
    # Print summary
    print()
    print("=" * 60)
    print("Summary")
    print("=" * 60)
    print(f"Total players processed: {stats['total_players']}")
    print(f"Slides created: {stats['slides_created']}")
    print(f"Images found and replaced: {stats['images_found']}")
    print(f"Images missing: {stats['images_missing']}")
    print(f"Errors: {stats['errors']}")
    print(f"Total slides in presentation: {len(prs.slides)}")
    print("=" * 60)


if __name__ == "__main__":
    main()
