#!/usr/bin/env python3
"""Quick script to verify images are present in the PowerPoint slides."""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path

pptx_path = Path("Re-Auction-2025-BCL-Players.pptx")

if not pptx_path.exists():
    print(f"Error: {pptx_path} not found")
    exit(1)

prs = Presentation(pptx_path)

print(f"Total slides: {len(prs.slides)}")
print("\nChecking first 10 slides for images:")
print("-" * 60)

for i, slide in enumerate(prs.slides[:10], 1):
    image_count = 0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_count += 1
            try:
                image = shape.image
                print(f"Slide {i}: Found image - Size: {len(image.blob)} bytes, Ext: {image.ext}")
            except:
                print(f"Slide {i}: Found image shape but couldn't read image data")
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub_shape in shape.shapes:
                if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_count += 1
                    try:
                        image = sub_shape.image
                        print(f"Slide {i}: Found image in group - Size: {len(image.blob)} bytes, Ext: {image.ext}")
                    except:
                        print(f"Slide {i}: Found image shape in group but couldn't read image data")
    
    if image_count == 0:
        print(f"Slide {i}: No images found")
    else:
        print(f"Slide {i}: Total images: {image_count}")

print("\n" + "=" * 60)
print("Checking a few slides from the middle:")
print("-" * 60)

# Check some slides from the middle
for i in [50, 100, 150]:
    if i < len(prs.slides):
        slide = prs.slides[i]
        image_count = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
        print(f"Slide {i+1}: {image_count} image(s) found")
