#!/usr/bin/env python3
"""
Sort slides in PowerPoint presentation by player name.
Keeps the first slide (template) in place and sorts all other slides alphabetically by player name.
"""

import sys
from pathlib import Path
from typing import List, Tuple
from pptx import Presentation


def extract_player_name_from_slide(slide) -> str:
    """Extract player name from the table on the slide."""
    for shape in slide.shapes:
        if hasattr(shape, 'has_table') and shape.has_table:
            table = shape.table
            # Player name is in the first cell (row 0, col 0)
            try:
                player_name = table.cell(0, 0).text.strip()
                return player_name if player_name else ''
            except:
                pass
    return ''


def sort_slides_by_name(pptx_path: Path) -> None:
    """Sort slides in the presentation by player name."""
    print("=" * 60)
    print("Sorting Slides by Player Name")
    print("=" * 60)
    print()
    
    # Load presentation
    print(f"Loading PowerPoint: {pptx_path.name}")
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"Error loading PowerPoint: {e}")
        sys.exit(1)
    
    if len(prs.slides) < 2:
        print("Not enough slides to sort (need at least 2 slides)")
        return
    
    print(f"Total slides: {len(prs.slides)}")
    print()
    
    # Extract player names from each slide
    print("Extracting player names from slides...")
    slide_data: List[Tuple[int, str, 'Slide']] = []
    
    # Keep template slide (first slide) separate
    template_slide = prs.slides[0]
    
    # Process all other slides
    for idx in range(1, len(prs.slides)):
        slide = prs.slides[idx]
        player_name = extract_player_name_from_slide(slide)
        slide_data.append((idx, player_name, slide))
        if player_name:
            print(f"  Slide {idx + 1}: {player_name}")
        else:
            print(f"  Slide {idx + 1}: (no name found)")
    
    print()
    
    # Sort by player name (case-insensitive)
    slide_data.sort(key=lambda x: x[1].lower() if x[1] else 'zzz')
    
    print("Sorted order:")
    for idx, (original_idx, name, _) in enumerate(slide_data, start=1):
        print(f"  Position {idx + 1}: {name if name else '(no name)'} (was slide {original_idx + 1})")
    print()
    
    # Reorder slides by manipulating the slide list
    print("Reordering slides...")
    
    # We need to reorder by manipulating the XML structure
    # Since python-pptx doesn't support direct reordering, we'll:
    # 1. Save slide data with their XML
    # 2. Remove all slides except template
    # 3. Re-add slides in sorted order
    
    import copy
    import tempfile
    
    # Save a backup first
    backup_path = pptx_path.with_suffix('.pptx.backup')
    prs.save(backup_path)
    print(f"Backup saved to: {backup_path.name}")
    
    # Store slide XML elements in sorted order
    sorted_slide_elements = []
    
    # Get template slide element (keep it first)
    template_element = template_slide._element
    
    # Get sorted slide elements
    for _, _, slide in slide_data:
        sorted_slide_elements.append(slide._element)
    
    # Remove all slides except template from the presentation
    # We need to work backwards to avoid index issues
    while len(prs.slides) > 1:
        rId = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[-1]
    
    # Now add slides in sorted order by copying their content
    for slide_element in sorted_slide_elements:
        # Get the original slide to copy from
        # We need to find it in the original presentation
        # Since we've already removed slides, we need to reload or use stored data
        pass
    
    # Better approach: Create new slides by copying from stored slide data
    # Reload the original to get slide content
    original_prs = Presentation(backup_path)
    
    # Create mapping of slide index to slide data
    slide_index_map = {idx: slide for idx, _, slide in slide_data}
    
    # Clear current slides except template
    while len(prs.slides) > 1:
        rId = prs.slides._sldIdLst[-1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[-1]
    
    # Add sorted slides
    for original_idx, _, _ in slide_data:
        original_slide = original_prs.slides[original_idx]
        
        # Create new slide with same layout
        new_slide = prs.slides.add_slide(original_slide.slide_layout)
        
        # Remove default shapes
        for shape in list(new_slide.shapes):
            sp = shape._element
            sp.getparent().remove(sp)
        
        # Copy all shapes from original slide
        for shape in original_slide.shapes:
            el = shape._element
            newel = copy.deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    
    # Save the sorted presentation
    print(f"Saving sorted presentation to: {pptx_path.name}")
    try:
        prs.save(pptx_path)
        print("✓ Presentation saved successfully")
    except Exception as e:
        print(f"✗ Error saving presentation: {e}")
        # Restore from backup
        import shutil
        shutil.copy(backup_path, pptx_path)
        print("Restored from backup")
        sys.exit(1)
    
    print()
    print("=" * 60)
    print("Summary")
    print("=" * 60)
    print(f"Template slide: Kept at position 1")
    print(f"Player slides sorted: {len(slide_data)}")
    print(f"Total slides: {len(prs.slides)}")
    print("=" * 60)


def main():
    """Main function."""
    script_dir = Path(__file__).parent
    pptx_path = script_dir / "Re-Auction-2025-BCL-Players.pptx"
    
    if not pptx_path.exists():
        print(f"Error: PowerPoint file not found: {pptx_path}")
        sys.exit(1)
    
    sort_slides_by_name(pptx_path)


if __name__ == "__main__":
    main()
