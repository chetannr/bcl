#!/usr/bin/env python3
"""Reset the PowerPoint to only have the template slide, then regenerate all player slides."""

from pptx import Presentation
from pathlib import Path
import sys

pptx_path = Path("Re-Auction-2025-BCL-Players.pptx")

if not pptx_path.exists():
    print(f"Error: {pptx_path} not found")
    sys.exit(1)

print("Loading presentation...")
prs = Presentation(pptx_path)

print(f"Current slides: {len(prs.slides)}")

# Keep only the first slide (template)
if len(prs.slides) > 1:
    print("Removing all slides except the template (slide 1)...")
    # Remove slides from the end to avoid index issues
    for i in range(len(prs.slides) - 1, 0, -1):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]
    
    print(f"Now has {len(prs.slides)} slide(s) (template only)")

# Save the reset presentation
prs.save(pptx_path)
print(f"Saved reset presentation to {pptx_path}")
print("\nNow run: python3 generate_player_slides.py")
